"""Document parser for PDF, Word, and PowerPoint files.

Extracts text content from document files and creates text fragments.
"""

from __future__ import annotations

import logging
from pathlib import Path

from src.core.models import FragmentType
from src.evidence.parsers.base import BaseParser, ParsedFragment, ParseResult

logger = logging.getLogger(__name__)


class DocumentParser(BaseParser):
    """Parser for document formats: PDF, DOCX, PPTX."""

    supported_formats = [".pdf", ".docx", ".pptx", ".doc", ".txt"]

    async def parse(self, file_path: str, file_name: str) -> ParseResult:
        """Parse a document file and extract text fragments.

        Routes to the appropriate sub-parser based on file extension.
        """
        ext = Path(file_name).suffix.lower()
        try:
            if ext == ".pdf":
                return await self._parse_pdf(file_path)
            elif ext in (".docx", ".doc"):
                return await self._parse_docx(file_path)
            elif ext == ".pptx":
                return await self._parse_pptx(file_path)
            elif ext == ".txt":
                return await self._parse_text(file_path)
            else:
                return ParseResult(error=f"Unsupported document format: {ext}")
        except Exception as e:  # Intentionally broad: parser library exceptions vary by format
            logger.exception("Failed to parse document: %s", file_name)
            return ParseResult(error=f"Parse error: {e}")

    async def _parse_pdf(self, file_path: str) -> ParseResult:
        """Extract text from PDF using pdfplumber."""
        import pdfplumber

        fragments: list[ParsedFragment] = []
        metadata: dict[str, str | int | float | bool | None] = {}

        with pdfplumber.open(file_path) as pdf:
            metadata["page_count"] = len(pdf.pages)
            for page_num, page in enumerate(pdf.pages, start=1):
                text = page.extract_text()
                if text and text.strip():
                    fragments.append(
                        ParsedFragment(
                            fragment_type=FragmentType.TEXT,
                            content=text.strip(),
                            metadata={"page": page_num},
                        )
                    )

                # Extract tables
                tables = page.extract_tables()
                for table_idx, table in enumerate(tables):
                    if table:
                        # Convert table to a readable string format
                        rows = []
                        for row in table:
                            cells = [str(cell) if cell is not None else "" for cell in row]
                            rows.append(" | ".join(cells))
                        table_text = "\n".join(rows)
                        if table_text.strip():
                            fragments.append(
                                ParsedFragment(
                                    fragment_type=FragmentType.TABLE,
                                    content=table_text,
                                    metadata={"page": page_num, "table_index": table_idx},
                                )
                            )

        return ParseResult(fragments=fragments, metadata=metadata)  # type: ignore[arg-type]

    async def _parse_docx(self, file_path: str) -> ParseResult:
        """Extract text from Word documents using python-docx."""
        from docx import Document

        fragments: list[ParsedFragment] = []
        metadata: dict[str, str | int | float | bool | None] = {}

        doc = Document(file_path)
        metadata["paragraph_count"] = len(doc.paragraphs)

        # Collect paragraphs into chunks
        current_text: list[str] = []
        paragraph_idx = 0

        for para in doc.paragraphs:
            if para.text.strip():
                current_text.append(para.text.strip())
            elif current_text:
                # Empty paragraph acts as section break
                fragments.append(
                    ParsedFragment(
                        fragment_type=FragmentType.TEXT,
                        content="\n".join(current_text),
                        metadata={"start_paragraph": paragraph_idx - len(current_text) + 1},
                    )
                )
                current_text = []
            paragraph_idx += 1

        # Flush remaining text
        if current_text:
            fragments.append(
                ParsedFragment(
                    fragment_type=FragmentType.TEXT,
                    content="\n".join(current_text),
                    metadata={"start_paragraph": paragraph_idx - len(current_text)},
                )
            )

        # Extract tables
        for table_idx, table in enumerate(doc.tables):
            rows = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                rows.append(" | ".join(cells))
            table_text = "\n".join(rows)
            if table_text.strip():
                fragments.append(
                    ParsedFragment(
                        fragment_type=FragmentType.TABLE,
                        content=table_text,
                        metadata={"table_index": table_idx},
                    )
                )

        return ParseResult(fragments=fragments, metadata=metadata)  # type: ignore[arg-type]

    async def _parse_pptx(self, file_path: str) -> ParseResult:
        """Extract text from PowerPoint presentations."""
        from pptx import Presentation

        fragments: list[ParsedFragment] = []
        metadata: dict[str, str | int | float | bool | None] = {}

        prs = Presentation(file_path)
        metadata["slide_count"] = len(prs.slides)

        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_text_parts: list[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text_parts.append(shape.text.strip())

            if slide_text_parts:
                fragments.append(
                    ParsedFragment(
                        fragment_type=FragmentType.TEXT,
                        content="\n".join(slide_text_parts),
                        metadata={"slide": slide_num},
                    )
                )

        return ParseResult(fragments=fragments, metadata=metadata)  # type: ignore[arg-type]

    async def _parse_text(self, file_path: str) -> ParseResult:
        """Extract content from plain text files."""
        with open(file_path, encoding="utf-8", errors="replace") as f:
            content = f.read()

        fragments: list[ParsedFragment] = []
        if content.strip():
            fragments.append(
                ParsedFragment(
                    fragment_type=FragmentType.TEXT,
                    content=content.strip(),
                    metadata={"char_count": len(content)},
                )
            )

        return ParseResult(
            fragments=fragments,
            metadata={"char_count": len(content)},
        )
